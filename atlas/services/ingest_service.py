"""
Ingest Service — extracts raw text content from any supported input type.

Supported:
  - .txt, .md           → read directly
  - .pdf                → PyMuPDF (no AVX2 needed)
  - .docx               → python-docx
  - .pptx               → python-pptx
  - .jpg/.png/.webp etc → Gemini 2.5 Flash vision
  - YouTube URL         → yt-dlp transcript/metadata
  - Instagram URL       → yt-dlp caption/metadata

Returns an ExtractedContent dataclass with raw text + metadata.
"""

import os
import re
import logging
import tempfile
from dataclasses import dataclass, field
from pathlib import Path

import httpx
from dotenv import load_dotenv

load_dotenv()
logger = logging.getLogger(__name__)

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")
GEMINI_URL = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent"

IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp", ".heic", ".heif"}
YOUTUBE_RE = re.compile(r"(https?://)?(www\.)?(youtube\.com/watch\?v=|youtu\.be/)[\w\-]+")
INSTAGRAM_RE = re.compile(r"(https?://)?(www\.)?instagram\.com/(p|reel|tv)/[\w\-]+")


@dataclass
class ExtractedContent:
    text: str
    source_type: str          # "document" | "image" | "video" | "url"
    filename: str = ""
    url: str = ""
    mime: str = ""
    chunks: list[str] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)


def _chunk_text(text: str, max_chars: int = 1500, overlap: int = 150) -> list[str]:
    """Split long text into overlapping chunks for embedding."""
    text = text.strip()
    if len(text) <= max_chars:
        return [text]
    chunks = []
    start = 0
    while start < len(text):
        end = start + max_chars
        # Try to break at a sentence boundary
        if end < len(text):
            break_at = text.rfind(". ", start, end)
            if break_at > start + max_chars // 2:
                end = break_at + 1
        chunks.append(text[start:end].strip())
        start = end - overlap
    return [c for c in chunks if c]


def _extract_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace").strip()


def _extract_pdf(path: Path) -> str:
    import fitz  # PyMuPDF
    doc = fitz.open(str(path))
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    return "\n\n".join(pages).strip()


def _extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip()).strip()


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        slides.append("\n".join(parts))
    return "\n\n".join(slides).strip()


VISION_PROMPT = (
    "Describe this image in detail. "
    "If there is any text visible, transcribe it exactly. "
    "Note any people, objects, locations, or important context. "
    "Be thorough — this will be used for semantic search."
)

OLLAMA_BASE = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")
OLLAMA_VISION_MODEL = os.getenv("OLLAMA_VISION_MODEL", "llava:7b")


async def _extract_image_ollama(path: Path) -> str:
    """Primary: use local Ollama llava vision model."""
    import base64
    b64 = base64.b64encode(path.read_bytes()).decode()
    payload = {
        "model": OLLAMA_VISION_MODEL,
        "prompt": VISION_PROMPT,
        "images": [b64],
        "stream": False,
    }
    async with httpx.AsyncClient(timeout=300) as client:
        r = await client.post(f"{OLLAMA_BASE}/api/generate", json=payload)
        r.raise_for_status()
        return r.json()["response"].strip()


async def _extract_image_gemini(path: Path) -> str:
    """Fallback: Gemini 2.0 Flash vision (cloud, free tier with daily quota)."""
    if not GEMINI_API_KEY:
        raise RuntimeError("GEMINI_API_KEY not set — cannot use Gemini fallback")
    import base64, asyncio as _asyncio
    ext = path.suffix.lower()
    mime_map = {
        ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
        ".png": "image/png",  ".webp": "image/webp",
        ".gif": "image/gif",  ".bmp": "image/bmp",
        ".heic": "image/heic",".heif": "image/heif",
    }
    mime = mime_map.get(ext, "image/jpeg")
    b64 = base64.b64encode(path.read_bytes()).decode()
    payload = {"contents": [{"parts": [
        {"text": VISION_PROMPT},
        {"inline_data": {"mime_type": mime, "data": b64}},
    ]}]}
    for attempt in range(3):
        async with httpx.AsyncClient(timeout=30) as client:
            r = await client.post(f"{GEMINI_URL}?key={GEMINI_API_KEY}", json=payload)
            if r.status_code == 429:
                wait = 20 * (attempt + 1)
                logger.warning(f"Gemini 429 — waiting {wait}s (attempt {attempt+1}/3)")
                await _asyncio.sleep(wait)
                continue
            r.raise_for_status()
            return r.json()["candidates"][0]["content"]["parts"][0]["text"].strip()
    raise RuntimeError("Gemini quota exceeded after 3 retries.")


_PRECISION_TRIGGERS = re.compile(
    r"bank|statement|transaction|receipt|invoice|balance|payment|salary|transfer|"
    r"sgd|usd|gbp|eur|myr|thb|amount|total|subtotal|tax|gst|vat|"
    r"account\s*no|acc\s*no|\d+\.\d{2}|"
    r"payslip|payroll|cheque|check|credit|debit|"
    r"document|contract|form|id\s*card|passport|license|certificate",
    re.IGNORECASE,
)


async def _extract_image(path: Path) -> str:
    """
    Smart vision routing:
    1. llava (local, fast) does a quick scan of the image.
    2. If the scan detects financial data, numbers, or text-heavy docs → escalate to Gemini for precision OCR.
    3. Otherwise return the llava result directly (saves Gemini API cost + latency).
    4. If llava is unavailable, go straight to Gemini.
    """
    llava_result = None

    try:
        logger.info(f"Vision fast-scan with {OLLAMA_VISION_MODEL}")
        llava_result = await _extract_image_ollama(path)
    except Exception as e:
        logger.warning(f"Ollama vision unavailable ({e}) — going straight to Gemini")
        return await _extract_image_gemini(path)

    # Check if llava detected anything that needs precision
    if _PRECISION_TRIGGERS.search(llava_result):
        logger.info("Precision content detected — escalating to Gemini for accurate OCR")
        try:
            return await _extract_image_gemini(path)
        except Exception as e:
            logger.warning(f"Gemini escalation failed ({e}) — using llava result")
            return llava_result

    logger.info("General image — using llava result")
    return llava_result


async def _extract_url(url: str) -> str:
    """Use yt-dlp to extract transcript/metadata from YouTube or Instagram URLs."""
    import yt_dlp

    ydl_opts = {
        "quiet": True,
        "no_warnings": True,
        "skip_download": True,
        "writesubtitles": False,
        "writeautomaticsub": True,
        "subtitleslangs": ["en"],
        "subtitlesformat": "vtt",
    }

    info = {}
    transcript_text = ""

    with tempfile.TemporaryDirectory() as tmpdir:
        ydl_opts["outtmpl"] = os.path.join(tmpdir, "%(id)s.%(ext)s")

        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=False) or {}

        # Try to get subtitle/transcript
        try:
            ydl_opts_subs = dict(ydl_opts)
            ydl_opts_subs["skip_download"] = False
            ydl_opts_subs["writeautomaticsub"] = True
            with yt_dlp.YoutubeDL(ydl_opts_subs) as ydl:
                ydl.extract_info(url, download=True)

            # Find .vtt file
            for f in Path(tmpdir).glob("*.vtt"):
                raw_vtt = f.read_text(encoding="utf-8", errors="replace")
                # Strip VTT markup, deduplicate lines
                lines = []
                seen = set()
                for line in raw_vtt.splitlines():
                    line = line.strip()
                    if not line or "-->" in line or line.startswith("WEBVTT") or line.isdigit():
                        continue
                    # Strip HTML tags
                    line = re.sub(r"<[^>]+>", "", line)
                    if line and line not in seen:
                        seen.add(line)
                        lines.append(line)
                transcript_text = " ".join(lines)
                break
        except Exception as e:
            logger.debug(f"Subtitle extraction failed (non-fatal): {e}")

    title = info.get("title", "")
    description = (info.get("description") or "")[:500]
    channel = info.get("uploader", "")
    duration = info.get("duration_string", "")

    parts = []
    if title:
        parts.append(f"Title: {title}")
    if channel:
        parts.append(f"Channel/Creator: {channel}")
    if duration:
        parts.append(f"Duration: {duration}")
    if description:
        parts.append(f"Description: {description}")
    if transcript_text:
        parts.append(f"Transcript:\n{transcript_text[:3000]}")

    return "\n\n".join(parts).strip()


async def extract(
    file_bytes: bytes | None = None,
    filename: str = "",
    url: str = "",
) -> ExtractedContent:
    """
    Main extraction entry point.
    Pass either file_bytes+filename OR a url.
    """
    # ── URL path ──────────────────────────────────────────────────────────────
    if url:
        is_yt = bool(YOUTUBE_RE.search(url))
        is_ig = bool(INSTAGRAM_RE.search(url))
        if is_yt or is_ig:
            source_label = "youtube" if is_yt else "instagram"
            text = await _extract_url(url)
            chunks = _chunk_text(text)
            return ExtractedContent(
                text=text, source_type="url", url=url,
                chunks=chunks, metadata={"platform": source_label}
            )
        raise ValueError(f"URL not supported: {url}. Supported: YouTube, Instagram.")

    # ── File path ─────────────────────────────────────────────────────────────
    if file_bytes is None:
        raise ValueError("Must provide file_bytes or url")

    ext = Path(filename).suffix.lower()

    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = Path(tmp.name)

    try:
        if ext in (".txt", ".md"):
            text = _extract_txt(tmp_path)
            source_type = "document"
        elif ext == ".pdf":
            text = _extract_pdf(tmp_path)
            source_type = "document"
        elif ext == ".docx":
            text = _extract_docx(tmp_path)
            source_type = "document"
        elif ext in (".pptx", ".ppt"):
            text = _extract_pptx(tmp_path)
            source_type = "document"
        elif ext in IMAGE_EXTS:
            text = await _extract_image(tmp_path)
            source_type = "image"
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    finally:
        tmp_path.unlink(missing_ok=True)

    chunks = _chunk_text(text)
    return ExtractedContent(
        text=text, source_type=source_type,
        filename=filename, chunks=chunks,
    )
